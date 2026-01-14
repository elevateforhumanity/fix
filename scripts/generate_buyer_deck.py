#!/usr/bin/env python3
"""
Generate Buyer-Ready PowerPoint Deck
Pulls content from repo docs, creates professional PPTX
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import os

# Helper to create RGB color
def rgb_color(r, g, b):
    """Return tuple for RGB color"""
    return (r, g, b)

# Output path
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), '../docs/Elevate_Licensing_Deck.pptx')

# Brand colors (as tuples)
BRAND_BLUE = rgb_color(30, 58, 95)  # #1e3a5f
BRAND_GREEN = rgb_color(22, 163, 74)  # #16a34a
DARK_TEXT = rgb_color(30, 30, 30)
LIGHT_TEXT = rgb_color(255, 255, 255)

def set_shape_color(shape, color_tuple):
    """Set shape fill color from RGB tuple"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = type('RGBColor', (), {'__init__': lambda s: None})()
    # Direct XML manipulation for color
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import qn
    spPr = shape._sp.spPr
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            srgbClr.set('val', '%02X%02X%02X' % color_tuple)

def set_font_color(font, color_tuple):
    """Set font color from RGB tuple - simplified approach"""
    # python-pptx handles this internally
    pass

def add_title_slide(prs, title, subtitle):
    """Add a title slide with dark background"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.theme_color = 1  # Dark theme
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, note=None):
    """Add a content slide with title and bullets"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BRAND_BLUE
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = LIGHT_TEXT
    
    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = bullet_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(12)
    
    # Note at bottom if provided
    if note:
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RgbColor(100, 100, 100)
    
    return slide

def add_diagram_slide(prs, title, diagram_text):
    """Add a slide with a text-based diagram"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BRAND_BLUE
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = LIGHT_TEXT
    
    # Diagram box
    diagram_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.8), Inches(8), Inches(4.5)
    )
    diagram_box.fill.solid()
    diagram_box.fill.fore_color.rgb = RgbColor(245, 245, 245)
    diagram_box.line.color.rgb = RgbColor(200, 200, 200)
    
    # Diagram text
    text_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(7), Inches(4))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = diagram_text
    p.font.size = Pt(16)
    p.font.name = "Consolas"
    p.font.color.rgb = DARK_TEXT
    
    return slide

def add_screenshot_placeholder_slide(prs, title, route, description):
    """Add a slide with screenshot placeholder"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BRAND_BLUE
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = LIGHT_TEXT
    
    # Screenshot placeholder
    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.5), Inches(8.5), Inches(4.5)
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RgbColor(230, 230, 230)
    placeholder.line.color.rgb = RgbColor(180, 180, 180)
    
    # Placeholder text
    ph_text = slide.shapes.add_textbox(Inches(2), Inches(3.2), Inches(6), Inches(1))
    tf = ph_text.text_frame
    p = tf.paragraphs[0]
    p.text = f"[Live Demo: {route}]"
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.8))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_TEXT
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a comparison table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BRAND_BLUE
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = LIGHT_TEXT
    
    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.75), Inches(1.5), Inches(8.5), Inches(4)).table
    
    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = LIGHT_TEXT
        p.font.size = Pt(14)
    
    # Rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_TEXT
    
    return slide

def add_cta_slide(prs, title, bullets, contact_info):
    """Add a call-to-action slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BRAND_BLUE
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = LIGHT_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(7), Inches(2.5))
    tf = bullet_box.text_frame
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{i+1}. {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = LIGHT_TEXT
        p.space_after = Pt(8)
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1.5))
    tf = contact_box.text_frame
    for i, line in enumerate(contact_info):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RgbColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def generate_deck():
    """Generate the complete buyer deck"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "Elevate for Humanity",
        "Enterprise Workforce Automation Platform"
    )
    
    # Slide 2: The Problem
    add_content_slide(
        prs,
        "The Problem: Manual Workflows Don't Scale",
        [
            "Paper-based enrollment and signatures",
            "No unified case record across parties",
            "Manual task assignment and tracking",
            "Audit trails scattered across systems",
            "Compliance reporting requires manual compilation"
        ]
    )
    
    # Slide 3: What the Platform Is
    add_content_slide(
        prs,
        "Automation-First Infrastructure",
        [
            "The platform serves as the system of record for programs",
            "Program activation is driven by required signatures",
            "Tasks, milestones, and reporting initialize automatically",
            "All actions logged for audit and reimbursement purposes"
        ],
        note="Credentials and instructional partners may be provided by licensee."
    )
    
    # Slide 4: System Architecture
    add_diagram_slide(
        prs,
        "System Architecture",
        """enrollment_cases (canonical record)
        ↓
apprentice_agreements (multi-party signatures)
        ↓
[DB Trigger: on_signature_added]
        ↓
checkSignatureCompleteness()
        ↓
case_tasks (auto-initialized from templates)
        ↓
case_events (append-only audit log)"""
    )
    
    # Slide 5: Signature-Driven Activation
    add_content_slide(
        prs,
        "Signature-Driven Activation",
        [
            "Student signature (enrollment consent)",
            "Employer signature (OJT commitment)",
            "Program holder signature (sponsorship agreement)",
            "System validates completeness automatically",
            "Case activates only when all required signatures received"
        ]
    )
    
    # Slide 6: Automated Task Initialization
    add_content_slide(
        prs,
        "Automated Task Initialization",
        [
            "Document upload tasks (student)",
            "Verification tasks (program holder)",
            "RAPIDS registration (program holder)",
            "Orientation completion (student)",
            "Employer agreement (employer)"
        ],
        note="10 tasks seeded for barber apprenticeship program."
    )
    
    # Slide 7: Audit & Compliance Ledger
    add_content_slide(
        prs,
        "Audit & Compliance Ledger",
        [
            "case_created - enrollment initiated",
            "signature_added - party signed agreement",
            "status_changed - case state transition",
            "tasks_initialized - auto-created from templates",
            "task_completed - milestone achieved",
            "document_uploaded - evidence attached"
        ],
        note="Append-only. Exportable for WIOA, RAPIDS, and agency reporting."
    )
    
    # Slide 8: Demo - Admin View
    add_screenshot_placeholder_slide(
        prs,
        "Demo: Admin View",
        "/demo/admin",
        "Case management dashboard • Signature status tracking • Task completion monitoring • Audit event timeline"
    )
    
    # Slide 9: Demo - Learner View
    add_screenshot_placeholder_slide(
        prs,
        "Demo: Learner View",
        "/demo/learner",
        "Personal case status • Required tasks and due dates • Document upload interface • Progress tracking"
    )
    
    # Slide 10: Demo - Employer View
    add_screenshot_placeholder_slide(
        prs,
        "Demo: Employer/Partner View",
        "/demo/employer",
        "Assigned cases • Signature requests • Verification tasks • Hours approval workflow"
    )
    
    # Slide 11: Licensing Model
    add_content_slide(
        prs,
        "Licensing Model",
        [
            "Platform access and automation infrastructure",
            "Case spine and workflow engine",
            "Audit logging and compliance reporting",
            "API access for integrations",
            "Implementation support and training"
        ],
        note="Pricing varies by scope, region, and automation requirements."
    )
    
    # Slide 12: Implementation Overview
    add_content_slide(
        prs,
        "Implementation Timeline",
        [
            "Week 1-2: Environment setup and configuration",
            "Week 3-4: Data migration and integration",
            "Week 5-6: User training and pilot",
            "Week 7+: Production launch and support"
        ],
        note="Typical implementation: 6-8 weeks."
    )
    
    # Slide 13: Why License vs Build
    add_table_slide(
        prs,
        "Why License vs Build",
        ["Build In-House", "License Platform"],
        [
            ["12-18 months", "6-8 weeks"],
            ["$500K-$1M+ dev cost", "Scope-based licensing"],
            ["Ongoing maintenance", "Managed updates"],
            ["Compliance risk", "Audit-ready"],
            ["No proven track record", "Production-tested"]
        ]
    )
    
    # Slide 14: Next Steps
    add_cta_slide(
        prs,
        "Next Steps",
        [
            "Live platform walkthrough (30 min)",
            "Scope discussion and requirements",
            "Licensing proposal",
            "Implementation planning"
        ],
        [
            "Phone: (317) 314-3757",
            "Email: elevate4humanityedu@gmail.com",
            "Schedule: elevateforhumanity.org/schedule"
        ]
    )
    
    # Save
    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"✓ Deck saved to: {OUTPUT_PATH}")
    
    return OUTPUT_PATH

if __name__ == "__main__":
    output = generate_deck()
    print(f"\nGenerated: {output}")
    print("\nSlides:")
    print("1. Title - Elevate for Humanity")
    print("2. The Problem - Manual Workflows Don't Scale")
    print("3. Automation-First Infrastructure")
    print("4. System Architecture")
    print("5. Signature-Driven Activation")
    print("6. Automated Task Initialization")
    print("7. Audit & Compliance Ledger")
    print("8. Demo: Admin View")
    print("9. Demo: Learner View")
    print("10. Demo: Employer/Partner View")
    print("11. Licensing Model")
    print("12. Implementation Timeline")
    print("13. Why License vs Build")
    print("14. Next Steps")
