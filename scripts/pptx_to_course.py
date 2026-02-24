"""
Convert PPTX slide decks into LMS-ready Markdown lessons + course manifest.
1 deck = 1 module, each non-empty slide = 1 lesson.

Usage:
  python3 scripts/pptx_to_course.py --in ./slides --out ./generated
"""

import argparse
import json
import os
import re
from dataclasses import dataclass, asdict
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Any, Optional

from pptx import Presentation
from tqdm import tqdm


def slugify(s: str) -> str:
    s = s.strip().lower()
    s = re.sub(r"[^\w\s-]", "", s)
    s = re.sub(r"[\s_-]+", "-", s)
    return s.strip("-") or "untitled"


def clean_text(s: str) -> str:
    s = s.replace("\u00a0", " ").replace("\r", "\n")
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


def shape_text(shape) -> str:
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return ""
    parts = []
    for p in shape.text_frame.paragraphs:
        runs = [r.text for r in p.runs] if p.runs else [p.text]
        line = "".join(runs).strip()
        if line:
            parts.append(line)
    return "\n".join(parts).strip()


def slide_all_text(slide) -> str:
    chunks = []
    for shape in slide.shapes:
        txt = shape_text(shape)
        if txt:
            chunks.append(txt)

    seen = set()
    uniq = []
    for c in chunks:
        c2 = clean_text(c)
        if c2 and c2 not in seen:
            uniq.append(c2)
            seen.add(c2)

    return clean_text("\n\n".join(uniq))


def slide_title_guess(slide, fallback: str) -> str:
    try:
        if slide.shapes.title and slide.shapes.title.text:
            t = clean_text(slide.shapes.title.text)
            if t:
                return t
    except Exception:
        pass

    txt = slide_all_text(slide)
    if txt:
        first = txt.splitlines()[0].strip()
        if len(first) <= 80:
            return first

    return fallback


@dataclass
class Lesson:
    id: str
    title: str
    source_pptx: str
    slide_number: int
    markdown_path: str


@dataclass
class Module:
    id: str
    title: str
    source_pptx: str
    lessons: List[str]
    markdown_path: str


@dataclass
class Course:
    id: str
    title: str
    created_at: str
    modules: List[str]


def write_text(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")


def build_slide_markdown(module_title: str, lesson_title: str, body: str, meta: Dict[str, Any]) -> str:
    frontmatter_lines = ["---"]
    for k, v in meta.items():
        frontmatter_lines.append(f"{k}: {json.dumps(v)}")
    frontmatter_lines.append("---")

    md = "\n".join(frontmatter_lines)
    md += f"\n\n# {lesson_title}\n\n"
    md += body.strip() + "\n"
    return md


def build_module_markdown(module_title: str, lesson_refs: List[Dict[str, Any]], meta: Dict[str, Any]) -> str:
    frontmatter_lines = ["---"]
    for k, v in meta.items():
        frontmatter_lines.append(f"{k}: {json.dumps(v)}")
    frontmatter_lines.append("---")

    md = "\n".join(frontmatter_lines)
    md += f"\n\n# {module_title}\n\n"
    md += "## Lessons\n\n"
    for lr in lesson_refs:
        md += f"- {lr['title']} (Slide {lr['slide_number']}) → `{lr['markdown_path']}`\n"
    md += "\n"
    return md


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--in", dest="input_dir", required=True, help="Folder containing .pptx files")
    ap.add_argument("--out", dest="output_dir", required=True, help="Output folder for generated content")
    ap.add_argument("--course-title", default="Qwik608 / EPA 608 Training Course", help="Course title")
    args = ap.parse_args()

    input_dir = Path(args.input_dir).resolve()
    output_dir = Path(args.output_dir).resolve()

    pptx_files = sorted(input_dir.glob("*.pptx"))
    if not pptx_files:
        raise SystemExit(f"No .pptx found in {input_dir}")

    lessons: Dict[str, Lesson] = {}
    modules: Dict[str, Module] = {}

    course_id = "qwik608-course"
    course = Course(
        id=course_id,
        title=args.course_title,
        created_at=datetime.utcnow().isoformat() + "Z",
        modules=[],
    )

    for pptx_path in tqdm(pptx_files, desc="Processing PPTX"):
        prs = Presentation(str(pptx_path))
        module_title = pptx_path.stem.replace("_", " ").strip()
        module_id = slugify(module_title)

        lesson_ids: List[str] = []
        lesson_refs_for_module_md: List[Dict[str, Any]] = []

        for i, slide in enumerate(prs.slides, start=1):
            body = slide_all_text(slide)
            if not body:
                continue

            lesson_title = slide_title_guess(slide, fallback=f"{module_title} — Slide {i}")
            lesson_id = f"{module_id}-s{i:02d}-{slugify(lesson_title)[:40]}"

            md_rel = Path("lessons") / f"{lesson_id}.md"
            md_abs = output_dir / md_rel

            meta = {
                "lesson_id": lesson_id,
                "module_id": module_id,
                "module_title": module_title,
                "source_pptx": pptx_path.name,
                "slide_number": i,
                "title": lesson_title,
            }

            md = build_slide_markdown(module_title, lesson_title, body, meta)
            write_text(md_abs, md)

            lesson = Lesson(
                id=lesson_id,
                title=lesson_title,
                source_pptx=pptx_path.name,
                slide_number=i,
                markdown_path=str(md_rel).replace("\\", "/"),
            )
            lessons[lesson_id] = lesson
            lesson_ids.append(lesson_id)

            lesson_refs_for_module_md.append(
                {
                    "id": lesson_id,
                    "title": lesson_title,
                    "slide_number": i,
                    "markdown_path": lesson.markdown_path,
                }
            )

        module_md_rel = Path("modules") / f"{module_id}.md"
        module_md_abs = output_dir / module_md_rel

        module_meta = {
            "module_id": module_id,
            "module_title": module_title,
            "source_pptx": pptx_path.name,
            "lesson_count": len(lesson_ids),
        }
        module_md = build_module_markdown(module_title, lesson_refs_for_module_md, module_meta)
        write_text(module_md_abs, module_md)

        module = Module(
            id=module_id,
            title=module_title,
            source_pptx=pptx_path.name,
            lessons=lesson_ids,
            markdown_path=str(module_md_rel).replace("\\", "/"),
        )
        modules[module_id] = module
        course.modules.append(module_id)

    manifest = {
        "course": asdict(course),
        "modules": {k: asdict(v) for k, v in modules.items()},
        "lessons": {k: asdict(v) for k, v in lessons.items()},
    }

    output_dir.mkdir(parents=True, exist_ok=True)
    (output_dir / "course.json").write_text(json.dumps(manifest, indent=2), encoding="utf-8")

    print(f"\nDone. {len(lessons)} lessons across {len(modules)} modules.")
    print(f"- Manifest: {output_dir / 'course.json'}")
    print(f"- Modules:  {output_dir / 'modules'}")
    print(f"- Lessons:  {output_dir / 'lessons'}")


if __name__ == "__main__":
    main()
