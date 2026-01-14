#!/usr/bin/env python3
"""
Generate Buyer-Ready PowerPoint Deck - Simplified Version
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_PATH = os.path.join(os.path.dirname(__file__), '../docs/Elevate_Licensing_Deck.pptx')

def add_slide(prs, title, content_lines, is_title_slide=False):
    """Add a slide with title and bullet content"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    if is_title_slide:
        # Centered title slide
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        if content_lines:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
            tf = sub_box.text_frame
            for i, line in enumerate(content_lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(20)
                p.alignment = PP_ALIGN.CENTER
    else:
        # Standard content slide
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        
        if content_lines:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
            tf = content_box.text_frame
            tf.word_wrap = True
            for i, line in enumerate(content_lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {line}" if not line.startswith(('•', '-', '[', '1', '2', '3', '4', '5')) else line
                p.font.size = Pt(18)
                p.space_after = Pt(8)
    
    return slide

def generate_deck():
    """Generate the complete buyer deck"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slides_data = [
        # Slide 1: Title
        {
            "title": "Elevate for Humanity",
            "content": [
                "Enterprise Workforce Automation Platform",
                "",
                "System of Record • Signature-Driven Activation",
                "Automated Tasks • Audit-Ready Reporting"
            ],
            "is_title": True
        },
        # Slide 2: The Problem
        {
            "title": "The Problem: Manual Workflows Don't Scale",
            "content": [
                "Paper-based enrollment and signatures",
                "No unified case record across parties",
                "Manual task assignment and tracking",
                "Audit trails scattered across systems",
                "Compliance reporting requires manual compilation"
            ]
        },
        # Slide 3: What the Platform Is
        {
            "title": "Automation-First Infrastructure",
            "content": [
                "The platform serves as the system of record for programs",
                "Program activation is driven by required signatures",
                "Tasks, milestones, and reporting initialize automatically",
                "All actions logged for audit and reimbursement purposes",
                "",
                "[Note: Credentials and instructional partners may be provided by licensee]"
            ]
        },
        # Slide 4: System Architecture
        {
            "title": "System Architecture",
            "content": [
                "enrollment_cases → Canonical case record",
                "    ↓",
                "apprentice_agreements → Multi-party signatures",
                "    ↓",
                "[DB Trigger] → checkSignatureCompleteness()",
                "    ↓",
                "case_tasks → Auto-initialized from templates",
                "    ↓",
                "case_events → Append-only audit log"
            ]
        },
        # Slide 5: Signature-Driven Activation
        {
            "title": "Signature-Driven Activation",
            "content": [
                "Student signature (enrollment consent)",
                "Employer signature (OJT commitment)",
                "Program holder signature (sponsorship agreement)",
                "System validates completeness automatically",
                "Case activates only when all required signatures received"
            ]
        },
        # Slide 6: Automated Task Initialization
        {
            "title": "Automated Task Initialization",
            "content": [
                "When case activates, system auto-creates:",
                "",
                "Document upload tasks (student)",
                "Verification tasks (program holder)",
                "RAPIDS registration (program holder)",
                "Orientation completion (student)",
                "Employer agreement (employer)",
                "",
                "[10 tasks seeded for barber apprenticeship program]"
            ]
        },
        # Slide 7: Audit & Compliance Ledger
        {
            "title": "Audit & Compliance Ledger",
            "content": [
                "Every action creates an audit record:",
                "",
                "case_created - enrollment initiated",
                "signature_added - party signed agreement",
                "status_changed - case state transition",
                "tasks_initialized - auto-created from templates",
                "task_completed - milestone achieved",
                "document_uploaded - evidence attached",
                "",
                "[Append-only. Exportable for WIOA, RAPIDS, agency reporting]"
            ]
        },
        # Slide 8: Demo - Admin View
        {
            "title": "Demo: Admin View",
            "content": [
                "[Live Demo: /demo/admin]",
                "",
                "Case management dashboard",
                "Signature status tracking",
                "Task completion monitoring",
                "Audit event timeline"
            ]
        },
        # Slide 9: Demo - Learner View
        {
            "title": "Demo: Learner View",
            "content": [
                "[Live Demo: /demo/learner]",
                "",
                "Personal case status",
                "Required tasks and due dates",
                "Document upload interface",
                "Progress tracking"
            ]
        },
        # Slide 10: Demo - Employer View
        {
            "title": "Demo: Employer/Partner View",
            "content": [
                "[Live Demo: /demo/employer]",
                "",
                "Assigned cases",
                "Signature requests",
                "Verification tasks",
                "Hours approval workflow"
            ]
        },
        # Slide 11: Licensing Model
        {
            "title": "Licensing Model",
            "content": [
                "Licensing covers:",
                "",
                "Platform access and automation infrastructure",
                "Case spine and workflow engine",
                "Audit logging and compliance reporting",
                "API access for integrations",
                "Implementation support and training",
                "",
                "[Pricing varies by scope, region, and automation requirements]"
            ]
        },
        # Slide 12: Implementation Timeline
        {
            "title": "Implementation Timeline",
            "content": [
                "Week 1-2: Environment setup and configuration",
                "Week 3-4: Data migration and integration",
                "Week 5-6: User training and pilot",
                "Week 7+: Production launch and support",
                "",
                "[Typical implementation: 6-8 weeks]"
            ]
        },
        # Slide 13: Why License vs Build
        {
            "title": "Why License vs Build",
            "content": [
                "BUILD IN-HOUSE          vs    LICENSE PLATFORM",
                "─────────────────────────────────────────────",
                "12-18 months                  6-8 weeks",
                "$500K-$1M+ dev cost           Scope-based licensing",
                "Ongoing maintenance           Managed updates",
                "Compliance risk               Audit-ready",
                "No proven track record        Production-tested"
            ]
        },
        # Slide 14: Next Steps
        {
            "title": "Next Steps",
            "content": [
                "1. Live platform walkthrough (30 min)",
                "2. Scope discussion and requirements",
                "3. Licensing proposal",
                "4. Implementation planning",
                "",
                "─────────────────────────────────────────────",
                "",
                "Phone: (317) 314-3757",
                "Email: elevate4humanityedu@gmail.com",
                "Schedule: elevateforhumanity.org/schedule"
            ],
            "is_title": True
        }
    ]
    
    for slide_data in slides_data:
        add_slide(
            prs,
            slide_data["title"],
            slide_data.get("content", []),
            slide_data.get("is_title", False)
        )
    
    # Save
    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"✓ Deck saved to: {OUTPUT_PATH}")
    
    return OUTPUT_PATH

if __name__ == "__main__":
    output = generate_deck()
    print(f"\nGenerated: {output}")
    print("\n=== Slides Generated ===")
    print("1. Title - Elevate for Humanity")
    print("2. The Problem - Manual Workflows Don't Scale")
    print("3. Automation-First Infrastructure")
    print("4. System Architecture")
    print("5. Signature-Driven Activation")
    print("6. Automated Task Initialization")
    print("7. Audit & Compliance Ledger")
    print("8. Demo: Admin View")
    print("9. Demo: Learner View")
    print("10. Demo: Employer/Partner View")
    print("11. Licensing Model")
    print("12. Implementation Timeline")
    print("13. Why License vs Build")
    print("14. Next Steps")
